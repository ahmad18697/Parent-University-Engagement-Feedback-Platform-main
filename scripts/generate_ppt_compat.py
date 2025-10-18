from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "Parent-University-Engagement-Presentation-compat.pptx"

slides = [
    ("Parent–University Engagement & Feedback Platform", "AI-enabled feedback intake, triage, and routing"),
    ("Problem & Motivation", "Fragmented channels; need centralized triage and routing"),
    ("Goals & Outcomes", "Portal + API; auto sentiment/category/priority; admin dashboard"),
    ("Architecture Overview", "FastAPI → Services → DB; UI templates; REST API"),
    ("Data Model", "Feedback fields + Department seeding"),
    ("AI Pipeline", "Rules + TextBlob; keywords + fuzzy; urgency cues"),
    ("Parallel Processing", "LangChain RunnableParallel for sentiment/category"),
    ("API Endpoints", "POST/GET feedback; GET departments; Swagger at /docs"),
    ("Web UI & Admin", "Form + results; filters; status; CSV export"),
    ("Demo Flow", "Submit → triage → filter → update status → export"),
    ("Config & Deployment", "Env vars; uvicorn; Docker"),
    ("Quality & Roadmap", "Tests; security; LLM routing; RBAC; analytics"),
    ("Summary", "Extensible, LLM-ready foundation with clear next steps"),
]

def build_simple():
    prs = Presentation()
    # Title slide
    title, subtitle = slides[0]
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

    # Content slides with simple bullets (no notes, no custom fonts)
    for title, bullets in slides[1:]:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        tf = s.shapes.placeholders[1].text_frame
        tf.clear()
        for i, bullet in enumerate(bullets.split("; ")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(20)

    prs.save(str(OUT_PATH))
    print(f"Saved presentation to: {OUT_PATH}")

if __name__ == "__main__":
    build_simple()
