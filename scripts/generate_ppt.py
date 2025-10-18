from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "Parent-University-Engagement-Presentation.pptx"

slides_data = [
    {
        "title": "Parent–University Engagement & Feedback Platform",
        "bullets": [
            "AI-powered intake, triage, and routing of parent feedback",
            "Built with FastAPI, SQLAlchemy, Jinja2; optional LangChain",
            "Clean web UI + REST API + admin dashboard",
        ],
        "notes": "Introduce the platform, purpose, and key technologies.",
    },
    {"title": "Problem & Motivation", "bullets": [
        "Fragmented channels slow response and reduce accountability",
        "Manual triage is error-prone; high-priority items are missed",
        "Need centralized intake, analytics, and automated routing",
    ], "notes": "Frame the pain points from institution + parent perspectives."},
    {"title": "Goals & Outcomes", "bullets": [
        "Single portal (web + API) to collect feedback",
        "Automatic sentiment, category, priority, and department mapping",
        "Admin dashboard with filters, status updates, CSV export",
    ], "notes": "Clarify what success looks like for stakeholders."},
    {"title": "Architecture Overview", "bullets": [
        "Request → FastAPI → AI Processing → Database → Response",
        "Modules: UI (templates/), API (routers/), Services (services/), ORM",
        "Storage: SQLite by default; easily swappable via DATABASE_URL",
    ], "notes": "High-level flow; AI is heuristics-first, LLM-ready."},
    {"title": "Data Model", "bullets": [
        "Feedback: sentiment, category, priority, department, status, created_at",
        "Department: seeded on startup for common units",
        "SQLAlchemy models under app/models.py",
    ], "notes": "Focus on fields relevant to triage and workflow."},
    {"title": "AI Pipeline (Heuristics-first)", "bullets": [
        "Sentiment: rules + optional TextBlob blend",
        "Categorization: keywords + fuzzy matching to map to department",
        "Priority: sentiment + keyword hits + urgency cues (e.g., 'urgent', 'unsafe')",
    ], "notes": "Balance reliability (rules) with optional ML for nuance."},
    {"title": "Parallel Processing (LangChain Runnables)", "bullets": [
        "RunnableParallel executes sentiment and category in parallel",
        "Combine node assigns final category/priority/department",
        "LLM routing can be enabled later via USE_LLM=true",
    ], "notes": "Explain why parallel execution improves responsiveness."},
    {"title": "API Endpoints", "bullets": [
        "POST /api/feedback – create feedback",
        "GET /api/feedback – list with filters",
        "GET /api/departments – list departments; Swagger at /docs",
    ], "notes": "Call out filter composability and Swagger for discovery."},
    {"title": "Web UI & Admin", "bullets": [
        "index.html: submission form shows resulting triage info",
        "feedback_list.html: filters, status badges, CSV export",
        "department_cards: per-department view of assignments",
    ], "notes": "Emphasize usability: quick triage, export, status updates."},
    {"title": "Demo Flow", "bullets": [
        "Submit via web form → triage results displayed",
        "Admin filter by department/sentiment/status",
        "Update status and export CSV",
    ], "notes": "Outline a 2-minute demo sequence."},
    {"title": "Configuration & Deployment", "bullets": [
        "Env: APP_NAME, DATABASE_URL, USE_LLM, OPENAI_API_KEY (if LLM)",
        "Local: uvicorn app.main:app --reload",
        "Docker: build/run with --env-file .env and -p 8000:8000",
    ], "notes": "DB can be swapped to Postgres via SQLAlchemy."},
    {"title": "Quality, Security & Roadmap", "bullets": [
        "Tests via pytest (services, routers, DB)",
        "PII handling, auth, rate limiting (future)",
        "Roadmap: pluggable LLM routing, RBAC, analytics, notifications",
    ], "notes": "Combine quality, privacy, and next steps succinctly."},
    {"title": "Summary", "bullets": [
        "Centralized, extensible, production-ready foundation",
        "Heuristics-first, LLM-ready design",
        "Clear next steps to scale and secure",
    ], "notes": "Close with the value proposition."},
]


def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_bullet_slide(prs, title_text, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)
    return slide


def add_notes(slide, notes_text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = notes_text


def build_presentation():
    prs = Presentation()
    add_title_slide(prs, "Parent–University Engagement & Feedback Platform", "AI-enabled feedback intake, triage, and routing")
    for s in slides_data:
        slide = add_bullet_slide(prs, s["title"], s["bullets"])
        if s.get("notes"):
            add_notes(slide, s["notes"]) 
    prs.save(str(OUT_PATH))
    print(f"Saved presentation to: {OUT_PATH}")


if __name__ == "__main__":
    build_presentation()
